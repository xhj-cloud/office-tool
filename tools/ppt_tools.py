"""
PowerPoint 演示文稿读写工具
支持 .pptx 格式的读取和生成
"""

import json
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from .json_repair import safe_parse_json


def _iter_shapes(shapes):
    """递归遍历所有形状（包括 group 组合形状内部的子形状）"""
    for shape in shapes:
        yield shape
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes)
        except Exception:
            pass


def _find_body_shape(slide):
    """查找正文内容占位符（BODY/OBJECT 类型）；找不到返回 None。

    注意：layout-0（标题页）的副标题占位符是 SUBTITLE 类型，不会被误认为正文框。
    """
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        try:
            ptype = shape.placeholder_format.type
        except Exception:
            continue
        if ptype in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return shape
    return None


def _add_body_textbox(slide, prs):
    """无正文占位符时（如空白 layout-6），添加默认位置的文本框承载内容"""
    left = int(prs.slide_width * 0.067)
    top = int(prs.slide_height * 0.24)
    width = int(prs.slide_width * 0.866)
    height = int(prs.slide_height * 0.64)
    return slide.shapes.add_textbox(left, top, width, height)


def read_pptx(file_path: str, mode: str = "full") -> str:
    """
    读取 PPT 文件内容

    Args:
        file_path: .pptx 文件路径
        mode: "full"(全部) / "outline"(大纲) / "notes"(含备注)

    Returns:
        JSON 格式的内容
    """
    try:
        prs = Presentation(file_path)
    except Exception as e:
        return json.dumps({"error": f"无法打开文件: {str(e)}"}, ensure_ascii=False)

    output = {
        "file": file_path,
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
    }

    slides_data = []
    for si, slide in enumerate(prs.slides):
        slide_info = {"slide_num": si + 1, "layout": slide.slide_layout.name}
        shapes = []

        for shape in _iter_shapes(slide.shapes):
            shape_info = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }

            if shape.has_text_frame:
                texts = []
                for para in shape.text_frame.paragraphs:
                    para_text = ""
                    for run in para.runs:
                        para_text += run.text
                    if para_text.strip():
                        texts.append(para_text)
                shape_info["texts"] = texts

            if shape.has_table:
                table = shape.table
                tdata = []
                for row in table.rows:
                    tdata.append([cell.text for cell in row.cells])
                shape_info["table"] = {"rows": len(table.rows), "cols": len(table.columns), "data": tdata}

            shapes.append(shape_info)

        slide_info["shapes"] = shapes

        if mode == "notes" and slide.has_notes_slide:
            slide_info["notes"] = slide.notes_slide.notes_text_frame.text

        slides_data.append(slide_info)

    output["slides"] = slides_data
    return json.dumps(output, ensure_ascii=False, indent=2)


def write_pptx(spec_json: str) -> str:
    """
    根据 JSON 规格生成 .pptx 文件

    spec_json 格式:
    {
        "output": "/path/to/output.pptx",
        "slide_width": 13.33,       // 英寸（16:9 默认）
        "slide_height": 7.5,
        "slides": [
            {
                "layout": 0,        // 0=标题, 1=标题+内容, 6=空白（默认）
                "title": "标题文字",
                "subtitle": "副标题",
                "bullets": ["要点1", "要点2", "要点3"],
                "table": {
                    "headers": ["列1", "列2"],
                    "rows": [["a", "b"]],
                    "left": 1.5, "top": 2.0, "width": 7.0, "height": 3.0
                }
            }
        ]
    }
    """
    spec, err = safe_parse_json(spec_json)
    if err:
        return json.dumps({"error": f"JSON 解析失败: {err}"}, ensure_ascii=False)

    output_path = spec.get("output")
    if not output_path:
        return json.dumps({"error": "必须指定 output 路径"}, ensure_ascii=False)

    prs = Presentation()

    # 页面尺寸（16:9）
    prs.slide_width = Inches(spec.get("slide_width", 13.33))
    prs.slide_height = Inches(spec.get("slide_height", 7.5))

    # 可用的 layouts: 0=Title, 1=Title and Content, 6=Blank
    layout_map = {
        0: 0,   # Title Slide
        1: 1,   # Title and Content
        6: 6,   # Blank
    }

    for slide_spec in spec.get("slides", []):
        layout_idx = slide_spec.get("layout", 6)
        layout_idx = layout_map.get(layout_idx, 6)
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # 标题
        if slide_spec.get("title") and slide.shapes.title:
            slide.shapes.title.text = slide_spec["title"]

        # 副标题
        if slide_spec.get("subtitle"):
            # 查找副标题占位符
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = slide_spec["subtitle"]
                    break

        # 要点（无正文占位符时自动添加文本框，不再静默丢弃）
        if slide_spec.get("bullets"):
            body_shape = _find_body_shape(slide) or _add_body_textbox(slide, prs)
            tf = body_shape.text_frame
            tf.clear()
            for i, bullet in enumerate(slide_spec["bullets"]):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

        # 表格
        if slide_spec.get("table"):
            tspec = slide_spec["table"]
            headers = tspec.get("headers", [])
            rows_data = tspec.get("rows", [])
            n_rows = len(rows_data) + 1
            n_cols = len(headers) or (len(rows_data[0]) if rows_data else 1)

            table_shape = slide.shapes.add_table(
                n_rows, n_cols,
                Inches(tspec.get("left", 1.5)),
                Inches(tspec.get("top", 2.0)),
                Inches(tspec.get("width", 7.0)),
                Inches(tspec.get("height", 3.0)),
            )
            table = table_shape.table

            # 表头
            for ci, h in enumerate(headers):
                cell = table.cell(0, ci)
                cell.text = h
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)

            # 数据
            for ri, row_data in enumerate(rows_data):
                for ci, val in enumerate(row_data):
                    cell = table.cell(ri + 1, ci)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return json.dumps({"success": True, "output": output_path, "slide_count": len(prs.slides)}, ensure_ascii=False)


def _replace_in_paragraph(para, find, replace):
    """在段落内替换所有 find 出现处，支持跨多个 run 的匹配。返回替换次数。

    PPT 文本常被格式差异/拼写检查切成多个 run，旧版只在单个 run 内匹配会静默漏掉。
    替换文本写入首个受影响 run（保留其格式），其余 run 中对应片段删除。
    """
    if not find:
        return 0
    count = 0
    while count < 1000:  # 安全上限：replace 包含 find 时防止无限循环
        full_text = "".join(run.text for run in para.runs)
        idx = full_text.find(find)
        if idx == -1:
            break
        end = idx + len(find)
        pos = 0
        affected = []
        for run in para.runs:
            rlen = len(run.text)
            r_start, r_end = pos, pos + rlen
            if r_end > idx and r_start < end:
                affected.append((run, max(idx - r_start, 0), min(end - r_start, rlen)))
            pos += rlen
        first_run, fs, fe = affected[0]
        first_run.text = first_run.text[:fs] + replace + first_run.text[fe:]
        for run, s, e in affected[1:]:
            run.text = run.text[:s] + run.text[e:]
        count += 1
    return count


def _replace_in_text_frame(tf, find, replace):
    """在文本框内查找替换文本（支持跨 run）。返回替换次数。"""
    total = 0
    for para in tf.paragraphs:
        total += _replace_in_paragraph(para, find, replace)
    return total


def edit_pptx(file_path: str, spec_json: str) -> str:
    """
    编辑已有 PowerPoint 文件（.pptx）

    spec_json 格式:
    {
        "output": "/path/to/output.pptx",   // 可选，不指定则覆盖原文件
        "operations": [
            {"op": "replace_all", "find": "旧", "replace": "新"},       // 全局替换所有形状/表格文本
            {"op": "set_slide_title", "slide": 1, "text": "新标题"},    // 修改指定页标题
            {"op": "add_slide", "layout": 6, "title": "标题", "bullets": ["要点1","要点2"]},  // 添加幻灯片
            {"op": "delete_slide", "slide": 2}                          // 删除指定页
        ]
    }
    """
    spec, err = safe_parse_json(spec_json)
    if err:
        return json.dumps({"error": f"JSON 解析失败: {err}"}, ensure_ascii=False)

    try:
        prs = Presentation(file_path)
    except Exception as e:
        return json.dumps({"error": f"无法打开文件: {str(e)}"}, ensure_ascii=False)

    replacements = 0
    skipped: list[str] = []

    for op in spec.get("operations", []):
        kind = op.get("op")
        try:
            if kind == "replace_all":
                find, replace = op.get("find", ""), op.get("replace", "")
                for slide in prs.slides:
                    for shape in _iter_shapes(slide.shapes):  # 含 group 内部形状
                        if shape.has_text_frame:
                            replacements += _replace_in_text_frame(shape.text_frame, find, replace)
                        if getattr(shape, "has_table", False) and shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    replacements += _replace_in_text_frame(cell.text_frame, find, replace)

            elif kind == "set_slide_title":
                idx = op.get("slide", 1) - 1
                if not (0 <= idx < len(prs.slides)):
                    skipped.append(f"set_slide_title: 第 {idx + 1} 页超出范围（共 {len(prs.slides)} 页）")
                elif prs.slides[idx].shapes.title is None:
                    skipped.append(f"set_slide_title: 第 {idx + 1} 页没有标题占位符，无法修改")
                else:
                    prs.slides[idx].shapes.title.text = op.get("text", "")

            elif kind == "add_slide":
                layout_idx = op.get("layout", 6)
                if layout_idx not in (0, 1, 6):
                    return json.dumps(
                        {"error": f"add_slide: layout 只支持 0（标题页）/ 1（标题+内容）/ 6（空白），收到 {layout_idx}"},
                        ensure_ascii=False)
                if layout_idx >= len(prs.slide_layouts):
                    return json.dumps(
                        {"error": f"add_slide: 该模板只有 {len(prs.slide_layouts)} 种版式，无法使用 layout {layout_idx}"},
                        ensure_ascii=False)
                slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                if op.get("title") and slide.shapes.title:
                    slide.shapes.title.text = op["title"]
                elif op.get("title"):
                    skipped.append(f"add_slide: 新页（layout {layout_idx}）没有标题占位符，标题未写入")
                if op.get("bullets"):
                    body = _find_body_shape(slide) or _add_body_textbox(slide, prs)
                    tf = body.text_frame
                    tf.clear()
                    for i, b in enumerate(op["bullets"]):
                        if i == 0:
                            tf.paragraphs[0].text = b
                        else:
                            p = tf.add_paragraph()
                            p.text = b

            elif kind == "delete_slide":
                idx = op.get("slide", 1) - 1
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                if 0 <= idx < len(slides):
                    xml_slides.remove(slides[idx])
                else:
                    skipped.append(f"delete_slide: 第 {idx + 1} 页超出范围（共 {len(slides)} 页）")

            else:
                return json.dumps({"error": f"未知操作: {kind}"}, ensure_ascii=False)

        except Exception as e:
            return json.dumps({"error": f"操作 '{kind}' 执行失败: {str(e)}"}, ensure_ascii=False)

    output_path = spec.get("output", file_path)
    prs.save(output_path)
    result = {"success": True, "output": output_path, "replacements": replacements}
    if skipped:
        result["skipped"] = skipped
    return json.dumps(result, ensure_ascii=False)
